from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
import time
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from backend.core.paths import metis_dir
from backend.runtime.document_converters import soffice_path


OFFICE_ARTIFACT_VALIDATION_SCHEMA = "metis.office_artifact_validation.v1"
OFFICE_ARTIFACT_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".pdf"}
RENDER_REQUIRED_EXTENSIONS = {".docx", ".pdf"}
MAX_RENDER_PAGES = 3


def is_office_artifact_path(path: str | Path) -> bool:
    return Path(str(path or "")).suffix.lower() in OFFICE_ARTIFACT_EXTENSIONS


def validate_office_artifact(path: str | Path, *, output_dir: str = "") -> dict[str, Any]:
    target = Path(str(path or "")).expanduser().resolve(strict=False)
    ext = target.suffix.lower()
    started = time.time()
    payload: dict[str, Any] = {
        "schema": OFFICE_ARTIFACT_VALIDATION_SCHEMA,
        "version": 1,
        "ok": False,
        "status": "failed",
        "path": str(target),
        "extension": ext,
        "checks": [],
        "inspect": {},
        "render": {},
        "evidence_paths": [],
        "duration_ms": 0,
    }
    try:
        if ext not in OFFICE_ARTIFACT_EXTENSIONS:
            return _finish(payload, started, "skipped", True, "not an office artifact")
        if not target.is_file():
            return _finish(payload, started, "failed", False, "file not found")
        if target.stat().st_size <= 0:
            return _finish(payload, started, "failed", False, "file is empty")

        evidence_dir = _validation_output_dir(target, output_dir)
        if ext == ".pdf":
            _validate_pdf(target, payload, evidence_dir)
        elif ext == ".docx":
            _validate_docx(target, payload, evidence_dir)
        elif ext == ".xlsx":
            _validate_xlsx(target, payload)
        elif ext == ".pptx":
            _validate_pptx(target, payload)

        checks = payload.get("checks") if isinstance(payload.get("checks"), list) else []
        ok = bool(checks) and all(bool(check.get("ok")) for check in checks if isinstance(check, dict))
        return _finish(payload, started, "validated" if ok else "failed", ok, _summary(payload))
    except Exception as exc:
        payload["error"] = f"{type(exc).__name__}: {exc}"
        return _finish(payload, started, "failed", False, str(payload["error"]))


def _validate_pdf(path: Path, payload: dict[str, Any], evidence_dir: Path) -> None:
    pages = 0
    metadata: dict[str, Any] = {}
    try:
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception:
            from PyPDF2 import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        pages = len(reader.pages)
        metadata = {str(k): str(v) for k, v in dict(reader.metadata or {}).items()}
        _add_check(payload, "inspect", pages > 0, f"{pages} page(s)")
    except Exception as exc:
        _add_check(payload, "inspect", False, f"{type(exc).__name__}: {exc}")
    payload["inspect"] = {"method": "pypdf", "pages": pages, "metadata": metadata}
    _render_pdf(path, payload, evidence_dir, start_page=1, end_page=min(max(pages, 1), MAX_RENDER_PAGES))


def _validate_docx(path: Path, payload: dict[str, Any], evidence_dir: Path) -> None:
    paragraphs = 0
    tables = 0
    headings: list[str] = []
    images = _count_docx_images(path)
    try:
        from docx import Document  # type: ignore

        doc = Document(str(path))
        paragraphs = len(doc.paragraphs)
        tables = len(doc.tables)
        headings = [
            paragraph.text
            for paragraph in doc.paragraphs
            if str(getattr(paragraph.style, "name", "") or "").lower().startswith("heading")
        ][:40]
        _add_check(payload, "inspect", paragraphs > 0 or tables > 0 or images > 0, f"{paragraphs} paragraph(s), {tables} table(s), {images} image(s)")
    except Exception as exc:
        _add_check(payload, "inspect", False, f"{type(exc).__name__}: {exc}")
    payload["inspect"] = {
        "method": "python-docx",
        "paragraphs": paragraphs,
        "tables": tables,
        "images": images,
        "headings": headings,
    }
    pdf_path = _docx_to_pdf(path, evidence_dir, payload)
    if pdf_path:
        _render_pdf(pdf_path, payload, evidence_dir, start_page=1, end_page=MAX_RENDER_PAGES)


def _validate_xlsx(path: Path, payload: dict[str, Any]) -> None:
    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets: list[dict[str, Any]] = []
        for ws in wb.worksheets:
            preview: list[list[Any]] = []
            for row in ws.iter_rows(max_row=10, values_only=True):
                preview.append([cell for cell in row])
            sheets.append({"name": ws.title, "max_row": ws.max_row, "max_column": ws.max_column, "preview": preview})
        _add_check(payload, "inspect", bool(sheets), f"{len(sheets)} sheet(s)")
        payload["inspect"] = {"method": "openpyxl", "sheets": sheets}
    except Exception as exc:
        payload["inspect"] = {"method": "openpyxl", "error": f"{type(exc).__name__}: {exc}"}
        _add_check(payload, "inspect", False, str(payload["inspect"]["error"]))


def _validate_pptx(path: Path, payload: dict[str, Any]) -> None:
    slides: list[dict[str, Any]] = []
    method = "python-pptx"
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        for index, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    texts.append(text)
            slides.append({"index": index, "texts": texts[:20]})
    except Exception as exc:
        method = "ooxml"
        slides = _pptx_ooxml_text(path)
        if not slides:
            payload["inspect"] = {"method": "python-pptx", "error": f"{type(exc).__name__}: {exc}"}
            _add_check(payload, "inspect", False, str(payload["inspect"]["error"]))
            return
    payload["inspect"] = {"method": method, "slides": len(slides), "slide_text": slides}
    _add_check(payload, "inspect", bool(slides), f"{len(slides)} slide(s)")


def _docx_to_pdf(path: Path, output_dir: Path, payload: dict[str, Any]) -> Path | None:
    soffice = soffice_path()
    if not soffice:
        _add_check(payload, "render", False, "LibreOffice soffice not found")
        payload["render"] = {"ok": False, "dependency": "LibreOffice soffice", "error": "soffice not found"}
        return None
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(path)]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120, check=False)
    pdf_path = output_dir / f"{path.stem}.pdf"
    if proc.returncode != 0 or not pdf_path.is_file():
        _add_check(payload, "render", False, (proc.stderr or proc.stdout or "DOCX to PDF conversion failed")[:500])
        payload["render"] = {
            "ok": False,
            "dependency": "LibreOffice soffice",
            "command": cmd,
            "stderr": (proc.stderr or "")[:2000],
            "stdout": (proc.stdout or "")[:2000],
        }
        return None
    evidence = payload.get("evidence_paths")
    if isinstance(evidence, list):
        evidence.append(str(pdf_path))
    return pdf_path


def _render_pdf(path: Path, payload: dict[str, Any], output_dir: Path, *, start_page: int, end_page: int) -> None:
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        _add_check(payload, "render", False, "Poppler pdftoppm not found")
        payload["render"] = {"ok": False, "dependency": "poppler pdftoppm", "error": "pdftoppm not found"}
        return
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = output_dir / path.stem
    first = max(1, int(start_page or 1))
    last = max(first, int(end_page or first))
    cmd = [pdftoppm, "-png", "-r", "150", "-f", str(first), "-l", str(last), str(path), str(prefix)]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120, check=False)
    images = sorted(str(item) for item in output_dir.glob(f"{path.stem}-*.png"))
    ok = proc.returncode == 0 and bool(images)
    _add_check(payload, "render", ok, f"{len(images)} rendered page image(s)" if images else (proc.stderr or "no rendered page image"))
    payload["render"] = {
        "ok": ok,
        "method": "pdftoppm",
        "images": images,
        "command": cmd,
        "stderr": (proc.stderr or "")[:2000],
    }
    evidence = payload.get("evidence_paths")
    if isinstance(evidence, list):
        evidence.extend(images)


def _validation_output_dir(path: Path, raw: str) -> Path:
    if raw:
        base = Path(raw).expanduser().resolve(strict=False)
    else:
        digest = hashlib.sha256(f"{path}|{path.stat().st_mtime_ns}".encode("utf-8", errors="replace")).hexdigest()[:12]
        base = metis_dir("artifacts", "office-validation", f"{path.stem}-{digest}")
    base.mkdir(parents=True, exist_ok=True)
    return base


def _add_check(payload: dict[str, Any], name: str, ok: bool, detail: str) -> None:
    checks = payload.setdefault("checks", [])
    if isinstance(checks, list):
        checks.append({"name": name, "ok": bool(ok), "detail": str(detail or "")[:500]})


def _finish(payload: dict[str, Any], started: float, status: str, ok: bool, summary: str) -> dict[str, Any]:
    payload["ok"] = bool(ok)
    payload["status"] = status
    payload["summary"] = summary
    payload["duration_ms"] = int((time.time() - started) * 1000)
    return payload


def _summary(payload: dict[str, Any]) -> str:
    checks = payload.get("checks") if isinstance(payload.get("checks"), list) else []
    failed = [str(item.get("detail") or item.get("name") or "") for item in checks if isinstance(item, dict) and not item.get("ok")]
    if failed:
        return "; ".join(item for item in failed if item)[:500] or "validation failed"
    if payload.get("extension") in RENDER_REQUIRED_EXTENSIONS:
        render = payload.get("render") if isinstance(payload.get("render"), dict) else {}
        return f"validated with render evidence ({len(render.get('images') or [])} image(s))"
    return "validated by inspection"


def _count_docx_images(path: Path) -> int:
    try:
        with zipfile.ZipFile(path) as archive:
            names = archive.namelist()
            media = [name for name in names if name.startswith("word/media/")]
            if media:
                return len(media)
            rels = archive.read("word/_rels/document.xml.rels")
            root = ElementTree.fromstring(rels)
            return sum(1 for elem in root.iter() if "image" in str(elem.attrib.get("Type", "")))
    except Exception:
        return 0


def _pptx_ooxml_text(path: Path) -> list[dict[str, Any]]:
    slides: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path) as archive:
            names = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
            for index, name in enumerate(names, start=1):
                root = ElementTree.fromstring(archive.read(name))
                texts = [str(elem.text or "").strip() for elem in root.iter() if elem.tag.endswith("}t") and str(elem.text or "").strip()]
                slides.append({"index": index, "texts": texts[:20]})
    except Exception:
        return []
    return slides


def validation_json(payload: dict[str, Any]) -> str:
    return json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True)
