from __future__ import annotations

import json
import zipfile
from pathlib import Path
from typing import Any, Dict, List
from xml.etree import ElementTree

from backend.tools.coding.foundation.core_mechanisms.path_security import (
    PathSecurityError,
    safe_path_for_read,
    safe_path_for_write,
)
from backend.runtime.office_artifact_validation import validate_office_artifact


def xlsx_create(
    output_path: str,
    title: str = "",
    rows: List[List[Any]] | str | None = None,
    sheets: List[Dict[str, Any]] | str | None = None,
) -> str:
    try:
        from openpyxl import Workbook  # type: ignore
        from openpyxl.styles import Font  # type: ignore
    except Exception as exc:
        return _json_error(
            "Missing XLSX dependency: openpyxl is required.",
            dependency="openpyxl",
            detail=f"{type(exc).__name__}: {exc}",
        )
    target = _resolve_output_path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    normalized_sheets = _normalize_sheets(rows=rows, sheets=sheets)
    if not normalized_sheets:
        normalized_sheets = [{"name": title or "Sheet1", "rows": []}]
    first = True
    summaries: List[Dict[str, Any]] = []
    for sheet in normalized_sheets:
        name = _safe_sheet_name(str(sheet.get("name") or title or "Sheet1"))
        ws = wb.active if first else wb.create_sheet(title=name)
        ws.title = name
        first = False
        if title and len(normalized_sheets) == 1:
            ws["A1"] = title
            ws["A1"].font = Font(bold=True, size=14)
            start_row = 3
        else:
            start_row = 1
        data_rows = _normalize_rows(sheet.get("rows"))
        for row_index, row in enumerate(data_rows, start=start_row):
            for col_index, value in enumerate(row, start=1):
                ws.cell(row=row_index, column=col_index, value=value)
        summaries.append({"name": ws.title, "rows": len(data_rows), "columns": max((len(row) for row in data_rows), default=0)})
    wb.save(str(target))
    validation = validate_office_artifact(target)
    return _json(
        {
            "ok": bool(validation.get("ok")),
            "status": "complete" if validation.get("ok") else "validation_failed",
            "output_path": str(target),
            "title": title,
            "sheets": summaries,
            "artifact_ready": bool(validation.get("ok")),
            "artifact_validation": validation,
            "error": "" if validation.get("ok") else str(validation.get("summary") or validation.get("error") or "artifact validation failed"),
        }
    )


def xlsx_inspect(path: str, max_rows: int = 20) -> str:
    source, error = _resolve_existing_file(path, expected_ext=".xlsx")
    if source is None:
        return _json_error(error or "XLSX file not found", path=path)
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:
        return _json_error(
            "Missing XLSX dependency: openpyxl is required.",
            path=str(source),
            dependency="openpyxl",
            detail=f"{type(exc).__name__}: {exc}",
        )
    wb = load_workbook(str(source), read_only=True, data_only=True)
    sheets: List[Dict[str, Any]] = []
    limit = max(1, min(int(max_rows or 20), 100))
    for ws in wb.worksheets:
        preview: List[List[Any]] = []
        for row in ws.iter_rows(max_row=limit, values_only=True):
            preview.append([cell for cell in row])
        sheets.append(
            {
                "name": ws.title,
                "max_row": ws.max_row,
                "max_column": ws.max_column,
                "preview": preview,
            }
        )
    return _json({"ok": True, "path": str(source), "sheets": sheets})


def pptx_create(
    output_path: str,
    title: str = "",
    slides: List[Dict[str, Any]] | str | None = None,
) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return _json_error(
            "Missing PPTX dependency: python-pptx is required.",
            dependency="python-pptx",
            detail=f"{type(exc).__name__}: {exc}",
        )
    target = _resolve_output_path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    normalized_slides = _normalize_slides(slides)
    if not normalized_slides:
        normalized_slides = [{"title": title or "Presentation", "bullets": []}]
    for index, item in enumerate(normalized_slides):
        slide_title = str(item.get("title") or (title if index == 0 else f"Slide {index + 1}") or "").strip()
        body = str(item.get("body") or item.get("text") or "").strip()
        bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
        layout = prs.slide_layouts[1] if body or bullets else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        placeholders = list(slide.placeholders)
        content = placeholders[1] if len(placeholders) > 1 else None
        if content is not None and hasattr(content, "text_frame"):
            frame = content.text_frame
            frame.clear()
            if body:
                frame.text = body
            for bullet in bullets:
                paragraph = frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
    prs.save(str(target))
    validation = validate_office_artifact(target)
    return _json(
        {
            "ok": bool(validation.get("ok")),
            "status": "complete" if validation.get("ok") else "validation_failed",
            "output_path": str(target),
            "title": title,
            "slides": len(normalized_slides),
            "artifact_ready": bool(validation.get("ok")),
            "artifact_validation": validation,
            "error": "" if validation.get("ok") else str(validation.get("summary") or validation.get("error") or "artifact validation failed"),
        }
    )


def pptx_inspect(path: str) -> str:
    source, error = _resolve_existing_file(path, expected_ext=".pptx")
    if source is None:
        return _json_error(error or "PPTX file not found", path=path)
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(source))
        slides = []
        for index, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    texts.append(text)
            slides.append({"index": index, "texts": texts[:20]})
        return _json({"ok": True, "path": str(source), "slides": len(prs.slides), "slide_text": slides})
    except Exception as exc:
        fallback = _pptx_ooxml_text(source)
        if fallback:
            return _json({"ok": True, "path": str(source), "slides": len(fallback), "slide_text": fallback, "python_pptx_error": f"{type(exc).__name__}: {exc}"})
        return _json_error(
            "Missing PPTX dependency: python-pptx is required.",
            path=str(source),
            dependency="python-pptx",
            detail=f"{type(exc).__name__}: {exc}",
        )


def _normalize_sheets(*, rows: List[List[Any]] | str | None, sheets: List[Dict[str, Any]] | str | None) -> List[Dict[str, Any]]:
    parsed_sheets = _parse_jsonish(sheets)
    if isinstance(parsed_sheets, list):
        return [item for item in parsed_sheets if isinstance(item, dict)]
    parsed_rows = _normalize_rows(rows)
    return [{"name": "Sheet1", "rows": parsed_rows}] if parsed_rows else []


def _normalize_rows(value: Any) -> List[List[Any]]:
    parsed = _parse_jsonish(value)
    if not isinstance(parsed, list):
        return []
    rows: List[List[Any]] = []
    for item in parsed:
        if isinstance(item, list):
            rows.append(item)
        else:
            rows.append([item])
    return rows


def _normalize_slides(value: List[Dict[str, Any]] | str | None) -> List[Dict[str, Any]]:
    parsed = _parse_jsonish(value)
    if isinstance(parsed, list):
        return [item for item in parsed if isinstance(item, dict)]
    if isinstance(parsed, dict):
        return [parsed]
    return []


def _parse_jsonish(value: Any) -> Any:
    if value is None or value == "":
        return None
    if isinstance(value, (list, dict)):
        return value
    try:
        return json.loads(str(value))
    except json.JSONDecodeError:
        return value


def _safe_sheet_name(value: str) -> str:
    cleaned = "".join("_" if ch in "[]:*?/\\'" else ch for ch in value).strip()
    return (cleaned or "Sheet1")[:31]


def _pptx_ooxml_text(path: Path) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path) as archive:
            names = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
            for index, name in enumerate(names, start=1):
                root = ElementTree.fromstring(archive.read(name))
                texts = [str(elem.text or "").strip() for elem in root.iter() if elem.tag.endswith("}t") and str(elem.text or "").strip()]
                slides.append({"index": index, "texts": texts[:20]})
    except Exception:
        return []
    return slides


def _resolve_existing_file(path: str, *, expected_ext: str) -> tuple[Path | None, str]:
    try:
        resolved = safe_path_for_read(str(path or ""))
    except PathSecurityError as exc:
        return None, str(exc)
    if not resolved.is_file():
        return None, f"{expected_ext.upper().lstrip('.')} file not found"
    if resolved.suffix.lower() != expected_ext:
        return None, f"Expected {expected_ext} file"
    return resolved, ""


def _resolve_output_path(path: str) -> Path:
    return safe_path_for_write(str(path or ""))


def _json(payload: Dict[str, Any]) -> str:
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _json_error(message: str, **extra: Any) -> str:
    payload = {"ok": False, "error": message}
    payload.update(extra)
    return _json(payload)
